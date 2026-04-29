"""
Document Extractor (Phase 6, Sprint 6.1).

Extracts unified content from non-IFC document formats:
- PDF (text-layer via pymupdf) -> markdown per page, drawing pages skipped
- DOCX (python-docx) -> markdown preserving headings, lists, tables
- XLSX (openpyxl) -> typed JSON {sheets: [{name, columns, rows, types}]}
- PPTX (python-pptx) -> markdown per slide

Returns a stateless ``DocumentExtractionResult`` describing one or more
``DocumentContentData`` payloads; persistence is Django's job (mirroring the
Phase 5 drawing-extractor pattern).

Drawing-vs-document classification for PDFs reuses
``drawing_extractor._looks_like_document_page`` so a mixed PDF can be split
between Phase 5 (drawings) and Phase 6 (documents) at dispatch time. This
extractor returns markdown ONLY for pages that look like documents; drawing
pages get an empty payload and ``is_document=False`` so the caller can route
them to the drawing extractor instead.

Failure-isolated: a bad page/slide/sheet logs a warning and the rest
continue. Catastrophic failures set ``success=False`` and ``error``.
"""

from __future__ import annotations

import os
import re
import time
from dataclasses import dataclass, field
from datetime import datetime
from typing import Any, Callable, Dict, List, Optional

from services.drawing_extractor import _looks_like_document_page


# ---------------------------------------------------------------------------
# Result types
# ---------------------------------------------------------------------------


@dataclass
class DocumentContentData:
    """
    One extracted document body. Whole-file formats (DOCX/XLSX/PPTX) emit a
    single payload; PDFs emit one payload per document page.
    """
    page_index: int = 0
    markdown_content: str = ""
    structured_data: Dict[str, Any] = field(default_factory=dict)
    page_count: int = 1
    structure: Dict[str, Any] = field(default_factory=dict)
    extracted_images: List[Dict[str, Any]] = field(default_factory=list)
    search_text: str = ""
    extraction_method: str = "structured"
    is_document: bool = True
    """False for PDF pages classified as drawings — the caller routes those to the drawing extractor."""


@dataclass
class DocumentExtractionResult:
    """Result of extracting one document source file."""
    success: bool = False
    format: str = ""  # 'pdf' | 'docx' | 'xlsx' | 'pptx'
    documents: List[DocumentContentData] = field(default_factory=list)
    log_entries: List[Dict[str, Any]] = field(default_factory=list)
    quality_report: Dict[str, Any] = field(default_factory=dict)
    duration_seconds: float = 0.0
    error: Optional[str] = None


LogFn = Callable[..., None]


# ---------------------------------------------------------------------------
# Public entry point
# ---------------------------------------------------------------------------


def extract_document(file_path: str, fmt: str) -> DocumentExtractionResult:
    """
    Extract document content from a file.

    Args:
        file_path: Local path to the source file.
        fmt: One of 'pdf', 'docx', 'xlsx', 'pptx' (lowercase).

    Returns:
        DocumentExtractionResult populated with one payload per page (PDF) or
        one payload total (DOCX/XLSX/PPTX).
    """
    result = DocumentExtractionResult(format=fmt)
    start = time.time()

    def log(level: str, stage: str, message: str, **details: Any) -> None:
        result.log_entries.append({
            "timestamp": datetime.now().isoformat(),
            "level": level,
            "stage": stage,
            "message": message,
            **details,
        })

    if not os.path.exists(file_path):
        result.error = f"File not found: {file_path}"
        log("error", "open", result.error)
        result.duration_seconds = time.time() - start
        return result

    try:
        if fmt == "pdf":
            _extract_pdf(file_path, result, log)
        elif fmt == "docx":
            _extract_docx(file_path, result, log)
        elif fmt == "xlsx":
            _extract_xlsx(file_path, result, log)
        elif fmt == "pptx":
            _extract_pptx(file_path, result, log)
        else:
            result.error = f"Unsupported format '{fmt}'"
            log("error", "open", result.error)
            result.duration_seconds = time.time() - start
            return result

        result.success = True
        document_pages = sum(1 for d in result.documents if d.is_document)
        drawing_pages = sum(1 for d in result.documents if not d.is_document)
        total_chars = sum(len(d.markdown_content) for d in result.documents)
        result.quality_report = {
            "document_count": len(result.documents),
            "document_pages": document_pages,
            "drawing_pages": drawing_pages,
            "total_chars": total_chars,
        }
        log(
            "info",
            "complete",
            f"Extracted {len(result.documents)} document(s) in {time.time() - start:.2f}s",
        )
    except Exception as exc:  # pragma: no cover - catastrophic
        result.success = False
        result.error = str(exc)
        log("error", "fatal", f"Extraction failed: {exc}", error=str(exc))

    result.duration_seconds = time.time() - start
    return result


# ---------------------------------------------------------------------------
# PDF extraction (pymupdf, text-layer only)
# ---------------------------------------------------------------------------


def _extract_pdf(file_path: str, result: DocumentExtractionResult, log: LogFn) -> None:
    """One payload per PDF page. Drawing-classified pages get is_document=False."""
    import fitz  # pymupdf

    pdf = fitz.open(file_path)
    try:
        page_count = pdf.page_count
        for page_index in range(page_count):
            try:
                page = pdf[page_index]
                rect = page.rect  # PDF points (1/72 inch)
                width_mm = float(rect.width) * 25.4 / 72.0
                height_mm = float(rect.height) * 25.4 / 72.0
                area_mm2 = max(width_mm * height_mm, 1.0)

                text = page.get_text("text") or ""
                text_density = len(text) / area_mm2
                is_doc = _looks_like_document_page(width_mm, height_mm, text_density)

                if not is_doc:
                    # Drawing page — record an empty payload so the caller can
                    # see what happened, but emit no markdown.
                    result.documents.append(DocumentContentData(
                        page_index=page_index,
                        markdown_content="",
                        page_count=page_count,
                        extraction_method="text_layer",
                        is_document=False,
                    ))
                    continue

                markdown = _pdf_page_to_markdown(page)
                headings = _extract_pdf_headings(page)
                result.documents.append(DocumentContentData(
                    page_index=page_index,
                    markdown_content=markdown,
                    page_count=page_count,
                    structure={"headings": headings},
                    search_text=_normalize_search_text(text),
                    extraction_method="text_layer",
                    is_document=True,
                ))
            except Exception as exc:
                log(
                    "warning",
                    "pdf_page",
                    f"Page {page_index} extraction failed: {exc}",
                    page_index=page_index,
                    error=str(exc),
                )
    finally:
        pdf.close()

    log(
        "info",
        "pdf",
        f"PDF: {len(result.documents)} page(s) processed",
        page_count=len(result.documents),
    )


def _pdf_page_to_markdown(page) -> str:
    """
    Lightly-formatted markdown from a PDF page.

    pymupdf gives us text blocks with bbox + font info; we use font size as a
    rough proxy for headings. Anything noticeably larger than median body
    size becomes an h2/h3, the rest stays as paragraphs.
    """
    blocks = page.get_text("dict")
    spans: List[Dict[str, Any]] = []
    for block in blocks.get("blocks", []):
        for line in block.get("lines", []):
            line_text_parts: List[str] = []
            line_size = 0.0
            for span in line.get("spans", []):
                txt = (span.get("text") or "").strip()
                if not txt:
                    continue
                line_text_parts.append(txt)
                line_size = max(line_size, float(span.get("size") or 0.0))
            joined = " ".join(line_text_parts).strip()
            if joined:
                spans.append({"text": joined, "size": line_size})
    if not spans:
        return (page.get_text("text") or "").strip()

    sizes = sorted(s["size"] for s in spans)
    median = sizes[len(sizes) // 2] if sizes else 0.0
    out: List[str] = []
    for s in spans:
        text = s["text"]
        size = s["size"]
        if median > 0 and size >= median * 1.5:
            out.append(f"## {text}")
        elif median > 0 and size >= median * 1.2:
            out.append(f"### {text}")
        else:
            out.append(text)
    return "\n\n".join(out).strip()


def _extract_pdf_headings(page) -> List[Dict[str, Any]]:
    """Pull a flat list of likely headings (oversized lines) for structure metadata."""
    blocks = page.get_text("dict")
    candidates: List[Dict[str, Any]] = []
    sizes: List[float] = []
    for block in blocks.get("blocks", []):
        for line in block.get("lines", []):
            for span in line.get("spans", []):
                size = float(span.get("size") or 0.0)
                txt = (span.get("text") or "").strip()
                if txt and size > 0:
                    sizes.append(size)
                    candidates.append({"text": txt, "size": size})
    if not sizes:
        return []
    median = sorted(sizes)[len(sizes) // 2]
    return [
        {"text": c["text"], "level": 2 if c["size"] >= median * 1.5 else 3}
        for c in candidates
        if c["size"] >= median * 1.2
    ]


# ---------------------------------------------------------------------------
# DOCX extraction (python-docx)
# ---------------------------------------------------------------------------


def _extract_docx(file_path: str, result: DocumentExtractionResult, log: LogFn) -> None:
    """Whole-file payload. Headings, paragraphs, lists, tables -> markdown."""
    import docx

    doc = docx.Document(file_path)
    md_parts: List[str] = []
    headings: List[Dict[str, Any]] = []
    tables_meta: List[Dict[str, Any]] = []

    body = doc.element.body
    # python-docx exposes paragraphs and tables on Document but we need
    # in-order iteration for fidelity. The body's child elements give us that.
    from docx.text.paragraph import Paragraph
    from docx.table import Table
    from docx.oxml.ns import qn

    for child in body.iterchildren():
        tag = child.tag
        if tag == qn('w:p'):
            para = Paragraph(child, doc)
            text = (para.text or "").strip()
            if not text:
                continue
            style_name = (para.style.name or "").strip() if para.style else ""
            if style_name.startswith("Heading"):
                level = _parse_heading_level(style_name)
                md_parts.append(f"{'#' * level} {text}")
                headings.append({"text": text, "level": level})
            elif style_name in ("List Bullet", "List Paragraph"):
                md_parts.append(f"- {text}")
            elif style_name in ("List Number",):
                md_parts.append(f"1. {text}")
            else:
                md_parts.append(text)
        elif tag == qn('w:tbl'):
            try:
                table = Table(child, doc)
                rendered, meta = _render_docx_table(table)
                if rendered:
                    md_parts.append(rendered)
                    tables_meta.append(meta)
            except Exception as exc:
                log("warning", "docx_table", f"Table render failed: {exc}", error=str(exc))

    markdown = "\n\n".join(md_parts).strip()
    plain_text = "\n".join(p for p in md_parts if p)
    result.documents.append(DocumentContentData(
        page_index=0,
        markdown_content=markdown,
        page_count=1,
        structure={"headings": headings, "tables": tables_meta},
        search_text=_normalize_search_text(plain_text),
        extraction_method="structured",
        is_document=True,
    ))
    log(
        "info",
        "docx",
        f"DOCX: {len(headings)} heading(s), {len(tables_meta)} table(s)",
        heading_count=len(headings),
        table_count=len(tables_meta),
    )


def _parse_heading_level(style_name: str) -> int:
    """'Heading 1' -> 1, 'Heading 2' -> 2, fallback to 2."""
    m = re.search(r"(\d+)", style_name)
    if not m:
        return 2
    try:
        return max(1, min(6, int(m.group(1))))
    except ValueError:
        return 2


def _render_docx_table(table) -> tuple[str, Dict[str, Any]]:
    """Render a python-docx Table as a markdown table; return (markdown, meta)."""
    rows: List[List[str]] = []
    for row in table.rows:
        rows.append([(cell.text or "").strip() for cell in row.cells])
    if not rows:
        return ("", {"row_count": 0, "col_count": 0})
    col_count = max(len(r) for r in rows)
    # Pad short rows
    rows = [r + [""] * (col_count - len(r)) for r in rows]
    header = rows[0]
    body = rows[1:] if len(rows) > 1 else []
    header_md = "| " + " | ".join(header) + " |"
    sep_md = "| " + " | ".join(["---"] * col_count) + " |"
    body_md = "\n".join("| " + " | ".join(r) + " |" for r in body)
    table_md = "\n".join(filter(None, [header_md, sep_md, body_md]))
    return (table_md, {"row_count": len(rows), "col_count": col_count})


# ---------------------------------------------------------------------------
# XLSX extraction (openpyxl)
# ---------------------------------------------------------------------------


def _extract_xlsx(file_path: str, result: DocumentExtractionResult, log: LogFn) -> None:
    """Whole-file payload. Each sheet -> typed JSON with column/row/type info."""
    import openpyxl
    from datetime import date, datetime as _dt

    wb = openpyxl.load_workbook(file_path, data_only=True, read_only=False)
    sheets_data: List[Dict[str, Any]] = []
    search_chunks: List[str] = []

    for sheet_name in wb.sheetnames:
        try:
            ws = wb[sheet_name]
            rows_iter = ws.iter_rows(values_only=True)
            rows: List[List[Any]] = []
            for raw_row in rows_iter:
                rows.append(list(raw_row))
            if not rows:
                sheets_data.append({
                    "name": sheet_name,
                    "columns": [],
                    "rows": [],
                    "types": [],
                })
                continue

            columns = [str(c) if c is not None else "" for c in rows[0]]
            body = rows[1:]
            col_count = len(columns)
            # Pad short rows so types-by-column has a defined position for each cell.
            body = [r + [None] * (col_count - len(r)) for r in body]
            types = _infer_column_types(body, col_count)

            json_rows: List[List[Any]] = []
            for row in body:
                json_row: List[Any] = []
                for cell in row:
                    if isinstance(cell, (_dt, date)):
                        json_row.append(cell.isoformat())
                    else:
                        json_row.append(cell)
                json_rows.append(json_row)

            sheets_data.append({
                "name": sheet_name,
                "columns": columns,
                "rows": json_rows,
                "types": types,
            })

            search_chunks.append(sheet_name)
            search_chunks.extend(columns)
            for row in body:
                for cell in row:
                    if cell is None:
                        continue
                    search_chunks.append(str(cell))
        except Exception as exc:
            log(
                "warning", "xlsx_sheet",
                f"Sheet '{sheet_name}' failed: {exc}",
                sheet_name=sheet_name, error=str(exc),
            )

    structured = {"sheets": sheets_data}
    search_text = _normalize_search_text(" ".join(search_chunks))
    result.documents.append(DocumentContentData(
        page_index=0,
        markdown_content="",
        structured_data=structured,
        page_count=len(sheets_data),
        structure={"sheet_names": [s["name"] for s in sheets_data]},
        search_text=search_text,
        extraction_method="structured",
        is_document=True,
    ))
    log(
        "info", "xlsx",
        f"XLSX: {len(sheets_data)} sheet(s)",
        sheet_count=len(sheets_data),
    )


def _infer_column_types(body: List[List[Any]], col_count: int) -> List[str]:
    """Per-column type inference: 'number' | 'date' | 'boolean' | 'string' | 'mixed' | 'empty'."""
    from datetime import date, datetime as _dt
    types: List[str] = []
    for col_idx in range(col_count):
        observed: set[str] = set()
        for row in body:
            cell = row[col_idx] if col_idx < len(row) else None
            if cell is None or (isinstance(cell, str) and cell.strip() == ""):
                continue
            if isinstance(cell, bool):
                observed.add("boolean")
            elif isinstance(cell, (int, float)):
                observed.add("number")
            elif isinstance(cell, (_dt, date)):
                observed.add("date")
            else:
                observed.add("string")
        if not observed:
            types.append("empty")
        elif len(observed) == 1:
            types.append(observed.pop())
        else:
            types.append("mixed")
    return types


# ---------------------------------------------------------------------------
# PPTX extraction (python-pptx)
# ---------------------------------------------------------------------------


def _extract_pptx(file_path: str, result: DocumentExtractionResult, log: LogFn) -> None:
    """Whole-file payload. Each slide -> markdown section with title + body bullets."""
    from pptx import Presentation

    prs = Presentation(file_path)
    sections: List[str] = []
    headings: List[Dict[str, Any]] = []
    search_chunks: List[str] = []

    for slide_index, slide in enumerate(prs.slides):
        try:
            title = _pptx_slide_title(slide)
            body_lines: List[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if title and shape == _pptx_title_shape(slide):
                    continue  # already handled
                for para in shape.text_frame.paragraphs:
                    text = (para.text or "").strip()
                    if not text:
                        continue
                    body_lines.append(f"- {text}")
                    search_chunks.append(text)

            heading_text = title or f"Slide {slide_index + 1}"
            sections.append(f"## {heading_text}")
            if body_lines:
                sections.extend(body_lines)
            headings.append({"text": heading_text, "level": 2, "slide": slide_index})
            if title:
                search_chunks.append(title)
        except Exception as exc:
            log(
                "warning", "pptx_slide",
                f"Slide {slide_index} failed: {exc}",
                slide_index=slide_index, error=str(exc),
            )

    markdown = "\n\n".join(sections).strip()
    result.documents.append(DocumentContentData(
        page_index=0,
        markdown_content=markdown,
        page_count=len(prs.slides),
        structure={"headings": headings},
        search_text=_normalize_search_text(" ".join(search_chunks)),
        extraction_method="structured",
        is_document=True,
    ))
    log(
        "info", "pptx",
        f"PPTX: {len(prs.slides)} slide(s)",
        slide_count=len(prs.slides),
    )


def _pptx_slide_title(slide) -> str:
    title_shape = _pptx_title_shape(slide)
    if title_shape is None:
        return ""
    text = (title_shape.text_frame.text or "").strip() if title_shape.has_text_frame else ""
    return text


def _pptx_title_shape(slide):
    try:
        return slide.shapes.title
    except Exception:
        return None


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------


def _normalize_search_text(text: str) -> str:
    """Lowercase + collapse whitespace for predictable LIKE/FTS behaviour."""
    if not text:
        return ""
    collapsed = re.sub(r"\s+", " ", text).strip()
    return collapsed.lower()
