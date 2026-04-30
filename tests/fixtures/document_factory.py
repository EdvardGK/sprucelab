"""
Programmatic PDF/DOCX/XLSX/PPTX fixtures for Phase 6 document-extractor tests.

Same shape as ``drawing_factory.py``: each builder takes ``out: Path``,
returns ``out``, lazy-imports its dependency. Generated files live in tmp
paths — nothing is committed to the repo.

Includes both English and Norwegian content because the Sprint 6.2 claim
extractor will need to handle ``skal``/``må``/``minimum``/``minst`` alongside
``shall``/``must``/``at least``. Building bilingual fixtures now means we
don't need to retrofit them when 6.2 lands.
"""
from __future__ import annotations

from pathlib import Path


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------


def build_pdf_text_document(out: Path) -> Path:
    """
    Single A4 portrait page with a heading and dense body text — char-density
    well above 0.05 chars/mm² so ``_looks_like_document_page`` classifies it
    as a document.

    The body deliberately mixes English and Norwegian normative language so
    the same fixture can drive 6.2 claim-extractor tests.

    Density math: A4 ≈ 62370 mm². 100 copies of a ~50-char sentence ≈ 5000
    chars → 0.08 chars/mm², comfortably above the 0.05 cutoff. This mirrors
    the Phase 5 ``build_pdf_a4_document`` density tuning.
    """
    import fitz  # pymupdf

    doc = fitz.open()
    # A4 portrait: 210 x 297 mm = 595 x 842 PDF points.
    page = doc.new_page(width=595, height=842)

    # Heading at top — the body text below carries the bulk of the chars.
    page.insert_text((50, 60), "Spesifikasjon REI60", fontsize=16)

    # 240-char sentence × 25 ≈ 6000 chars / 62370 mm² ≈ 0.096 chars/mm²,
    # comfortably above the 0.05 cutoff. pymupdf's insert_textbox refuses
    # to write anything if the text doesn't fit — keep the count below ~30.
    body = (
        "Vegger i branncelle skal vaere REI60. "
        "Walls in fire compartments must be REI60. "
        "Minimum ventilasjon er 7.0 l/s per person. "
        "At least 7.0 L/s per person of fresh air shall be provided. "
        "Lorem ipsum dolor sit amet, consectetur adipiscing elit. "
    ) * 25
    rect = fitz.Rect(40, 80, 555, 800)
    page.insert_textbox(rect, body.strip(), fontsize=10, lineheight=1.2)
    doc.save(str(out))
    doc.close()
    return out


def build_pdf_multipage_doc_and_drawing(out: Path) -> Path:
    """
    Two-page PDF: page 0 is a dense A4 document, page 1 is a sparse A3
    drawing. The PDF dispatch path runs both extractors over the file —
    document extractor produces one DocumentContent for page 0 and skips
    page 1; drawing extractor produces two DrawingSheet rows (one per page).
    """
    import fitz

    doc = fitz.open()
    # Page 0: A4 document — same density tuning as build_pdf_text_document.
    page0 = doc.new_page(width=595, height=842)
    page0.insert_text((50, 60), "Project Specification", fontsize=16)
    # ~100-char sentence × 50 ≈ 5000 chars on A4 → density well above 0.05.
    body = (
        "All ventilation systems shall meet TEK17 minimum requirements. "
        "Vegger skal ha brannmotstand REI60. "
    ) * 50
    rect0 = fitz.Rect(40, 80, 555, 800)
    page0.insert_textbox(rect0, body.strip(), fontsize=10, lineheight=1.2)

    # Page 1: A3 landscape drawing (sparse so it stays a drawing).
    page1 = doc.new_page(width=1190, height=842)
    page1.insert_text((400, 50), "DRAWING TITLE", fontsize=14)

    doc.save(str(out))
    doc.close()
    return out


def build_pdf_claim_corpus(out: Path) -> Path:
    """
    Single A4 PDF densely seeded with NB+EN normative statements.

    Used by the Sprint 6.2 claim-extractor unit and e2e tests. Statements
    cover the four predicate patterns the heuristic extractor recognizes:
    fire-resistance class, ventilation flow rate, U-value, acoustic dB.
    A descriptive past-tense line is included so the extractor's reject
    branch is exercised in the e2e test, not just the unit test.

    Density-tuned to clear the 0.05 chars/mm² document cutoff (see Phase 5
    fixture math) — sentence-of-claims block × 6 plus a body filler.
    """
    import fitz  # pymupdf

    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    page.insert_text((50, 60), "Project Specification", fontsize=16)

    claims_block = (
        "Vegger i branncelle skal vaere REI60. "
        "External walls must be REI30. "
        "Minimum 7 l/s per person fresh air shall be provided. "
        "U-value shall not exceed 0.18 W/m2K for external walls. "
        "Internal walls shall be at least 50 dB acoustic class. "
        "Doors must be EI30. "
        "Doerer skal vaere EI30. "
        "External walls were REI60 historically. "
    )
    body = (claims_block * 6) + (
        "Lorem ipsum dolor sit amet, consectetur adipiscing elit. "
    ) * 30
    rect = fitz.Rect(40, 80, 555, 800)
    page.insert_textbox(rect, body.strip(), fontsize=10, lineheight=1.2)
    doc.save(str(out))
    doc.close()
    return out


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------


def build_docx_with_headings(out: Path) -> Path:
    """
    DOCX with H1/H2 headings, a bulleted list, and a small table.

    Exercises the markdown rendering paths in ``_extract_docx``: heading
    detection, list formatting, table row joining.
    """
    import docx

    doc = docx.Document()
    doc.add_heading("Project Specification", level=1)
    doc.add_paragraph("This is the introductory paragraph.")
    doc.add_heading("Fire Resistance", level=2)
    doc.add_paragraph(
        "Walls in fire compartments shall be REI60. "
        "Vegger i branncelle skal vaere REI60."
    )
    doc.add_heading("Materials", level=2)
    doc.add_paragraph("Acceptable materials:", style=None)
    doc.add_paragraph("Concrete", style="List Bullet")
    doc.add_paragraph("Steel", style="List Bullet")
    doc.add_paragraph("Timber (CLT)", style="List Bullet")

    table = doc.add_table(rows=3, cols=3)
    table.rows[0].cells[0].text = "Element"
    table.rows[0].cells[1].text = "Class"
    table.rows[0].cells[2].text = "Notes"
    table.rows[1].cells[0].text = "Wall"
    table.rows[1].cells[1].text = "REI60"
    table.rows[1].cells[2].text = "Fire compartment"
    table.rows[2].cells[0].text = "Door"
    table.rows[2].cells[1].text = "EI30"
    table.rows[2].cells[2].text = "Egress"

    doc.save(str(out))
    return out


# ---------------------------------------------------------------------------
# XLSX
# ---------------------------------------------------------------------------


def build_xlsx_with_typed_data(out: Path) -> Path:
    """
    Two-sheet XLSX: a 'Quantities' sheet with mixed numeric + string columns
    and a 'Schedule' sheet with date columns. The extractor must preserve
    types (numbers stay numbers, dates serialize as ISO strings).
    """
    import openpyxl
    from datetime import date

    wb = openpyxl.Workbook()
    ws1 = wb.active
    ws1.title = "Quantities"
    ws1.append(["Element", "Count", "Area_m2", "Approved"])
    ws1.append(["Wall", 42, 350.5, True])
    ws1.append(["Slab", 12, 800.0, True])
    ws1.append(["Door", 18, 0.0, False])

    ws2 = wb.create_sheet("Schedule")
    ws2.append(["Phase", "Start", "End"])
    ws2.append(["Foundation", date(2026, 5, 1), date(2026, 6, 15)])
    ws2.append(["Frame", date(2026, 6, 16), date(2026, 9, 30)])

    wb.save(str(out))
    return out


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------


def build_pptx_with_slides(out: Path) -> Path:
    """
    Three-slide PPTX with a title + bulleted body on each slide. Exercises
    the per-slide markdown emission and the title-shape detection.
    """
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank_layout = prs.slide_layouts[1]  # title + content

    titles_and_bullets = [
        ("Project Brief", ["Goal: deliver MVP", "Owner: BIM Coordinator"]),
        ("Brannkrav", ["Vegger skal vaere REI60", "Doerer skal vaere EI30"]),
        ("Schedule", ["Phase 1: Foundation", "Phase 2: Frame"]),
    ]

    for title, bullets in titles_and_bullets:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.title.text = title
        # First placeholder besides title is the body text frame
        body = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx != 0:  # idx 0 is the title
                body = ph
                break
        if body is None:
            txbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
            body = txbox
        tf = body.text_frame
        tf.text = bullets[0]
        for extra in bullets[1:]:
            p = tf.add_paragraph()
            p.text = extra

    prs.save(str(out))
    return out
